#!/usr/bin/env python3
"""
业务快速上手演示文稿 — 内部分享
主题：如何快速了解一个陌生业务
时长：30分钟
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════
# Design System
# ═══════════════════════════════════════════════════════════

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
OLIVE       = RGBColor(0x5e, 0x5d, 0x59)
STONE       = RGBColor(0x87, 0x86, 0x7f)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)
TAG_BG      = RGBColor(0xee, 0xf2, 0xf7)

SERIF = "STSong"
SANS  = "PingFang SC"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(bg=PARCHMENT):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s

def txt(slide, text, l, t, w, h, font=SANS, size=18, bold=False,
        color=NEAR_BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '1800'})
    lnSpc.append(spcPct); pPr.append(lnSpc)
    spcAft = pPr.makeelement(qn('a:spcAft'), {'val': '200000'})
    pPr.append(spcAft)
    run = p.add_run(); run.text = text; run.font.name = font
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return tb

def multi_txt(slide, lines, l, t, w, h, font=SANS, size=16, color=DARK_WARM, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        pPr = p._p.get_or_add_pPr()
        # 行距 180%，段后间距 20pt，段落间距更宽松
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '1800'})
        lnSpc.append(spcPct); pPr.append(lnSpc)
        spcAft = pPr.makeelement(qn('a:spcAft'), {'val': '200000'})
        pPr.append(spcAft)
        run = p.add_run()
        run.text = item.get('text', ''); run.font.name = item.get('font', font)
        run.font.size = Pt(item.get('size', size)); run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', color)
    return tb

def card(slide, l, t, w, h, fill=IVORY, border=BORDER, bw=0.5):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(bw)
    c.shadow.inherit = False
    return c

IN = lambda inches: Inches(inches)

def page_num(slide, n, total):
    txt(slide, f"{n} / {total}", Inches(12), Inches(6.85), Inches(1.2), Inches(0.35),
        font=SANS, size=11, color=STONE, align=PP_ALIGN.RIGHT)

def brand_bar(slide, l, t, w, h=Inches(0.06)):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    b.fill.solid(); b.fill.fore_color.rgb = BRAND
    b.line.fill.background(); b.shadow.inherit = False
    return b

def page_header(slide, eyebrow, title, bar_y=Inches(1.6)):
    """页眉：eyebrow + 标题 + 分隔线"""
    txt(slide, eyebrow, IN(1), IN(0.4), IN(10), IN(0.35),
        font=SANS, size=12, color=STONE)
    txt(slide, title, IN(1), IN(0.82), IN(11.33), IN(0.65),
        font=SERIF, size=30, color=NEAR_BLACK)
    brand_bar(slide, IN(1), bar_y, IN(11.33))

def chapter_divider(slide, number, title):
    s = slide
    txt(s, number, IN(0.8), IN(0.4), IN(2), IN(0.8),
        font=SERIF, size=26, color=WHITE)
    txt(s, title, IN(1), IN(3), IN(11.33), IN(1.5),
        font=SERIF, size=48, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# Slide 1: Cover
# ═══════════════════════════════════════════════════════════

s = blank_slide()
txt(s, "如何快速了解一个陌生业务", IN(1), IN(2.5), IN(11.33), IN(1.5),
    font=SERIF, size=44, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
brand_bar(s, IN(6.17), IN(4.1), IN(1))
txt(s, "一套可复用的四步框架", IN(1), IN(4.4), IN(11.33), IN(0.8),
    font=SANS, size=20, color=OLIVE, align=PP_ALIGN.CENTER)
txt(s, "内部分享  ·  2026.04", IN(1), IN(6.2), IN(11.33), IN(0.4),
    font=SANS, size=13, color=STONE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# Slide 2: Opening
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "为什么分享这个", "从陌生到上手，最快的人是怎么做到的")
multi_txt(s, [
    {'text': '他在商业化方向上手极快——换一条业务线，一个月内就能给出清晰的策略判断。', 'size': 20},
    {'text': '', 'size': 12},
    {'text': '核心不是经验积累，而是有一套结构化的方法。', 'size': 20},
    {'text': '', 'size': 12},
    {'text': '今天把这套方法拆成四步，任何人都可以直接复用。', 'size': 20, 'bold': True, 'color': BRAND},
], IN(1), IN(2.0), IN(11.33), IN(4.2))
page_num(s, 2, 19)

# ═══════════════════════════════════════════════════════════
# Slide 3: TOC
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "目录", "四步法总览")

steps = [
    ("01", "搭建业务的「骨架模型」", "商业模式画布简化版，9个问题装满一张纸"),
    ("02", "用穿透式问题快速深挖", "四个维度，每组灵魂拷问逼迫你进入本质"),
    ("03", "三套调研动作交叉验证", "访谈、看数据、代入角色，互相核对"),
    ("04", "输出一页纸健康度诊断", "不写长报告，只填一张表加三个结论"),
]
for i, (num, title, desc) in enumerate(steps):
    y = IN(2.0 + i * 1.15)
    txt(s, num, IN(1.2), y, IN(0.7), IN(0.55),
        font=SERIF, size=26, color=BRAND)
    txt(s, title, IN(2.2), y, IN(5), IN(0.5),
        font=SERIF, size=22, color=NEAR_BLACK)
    txt(s, desc, IN(2.2), y + IN(0.5), IN(10), IN(0.35),
        font=SANS, size=13, color=OLIVE)

page_num(s, 3, 19)

# ═══════════════════════════════════════════════════════════
# Slide 4: Chapter 1 Divider
# ═══════════════════════════════════════════════════════════

s = blank_slide(BRAND)
chapter_divider(s, "01", "搭建业务的「骨架模型」")

# ═══════════════════════════════════════════════════════════
# Slide 5: BMC 9 Questions
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第一步 · 骨架模型", "商业模式画布 9 问 — 填不上的就是关键信息缺口")

questions = [
    ("1", "为谁", "目标客户细分"),
    ("2", "解决什么", "价值主张 · 痛点/爽点"),
    ("3", "怎么触达", "渠道通路 · 营销和销售"),
    ("4", "怎么维护", "客户关系 · 订阅/社群/1对1"),
    ("5", "靠什么赚钱", "收入来源 · 定价和收费对象"),
    ("6", "要做什么事", "关键业务活动 · 生产/研发/运营"),
    ("7", "有什么资源", "核心资源 · 人/技术/牌照/品牌/数据"),
    ("8", "和谁合作", "重要伙伴 · 供应商/分销商/平台"),
    ("9", "成本花在哪", "成本结构 · 固定 vs 变动"),
]

cw = IN(3.5); ch = IN(1.2)
gx = IN(0.35); gy = IN(0.2)
sx = IN(1.1); sy = IN(1.9)

for i, (num, q, a) in enumerate(questions):
    row, col = divmod(i, 3)
    x = sx + col * (cw + gx)
    y = sy + row * (ch + gy)
    card(s, x, y, cw, ch)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(x + IN(0.12)), Emu(y + IN(0.12)),
                               Emu(IN(0.36)), Emu(IN(0.36)))
    badge.fill.solid(); badge.fill.fore_color.rgb = BRAND
    badge.line.fill.background(); badge.shadow.inherit = False
    txt(s, num, Emu(x + IN(0.12)), Emu(y + IN(0.12)),
        IN(0.36), IN(0.36), font=SERIF, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, q, x + IN(0.6), y + IN(0.12), cw - IN(0.75), IN(0.4),
        font=SERIF, size=17, color=NEAR_BLACK)
    txt(s, a, x + IN(0.6), y + IN(0.6), cw - IN(0.75), IN(0.48),
        font=SANS, size=12, color=OLIVE)

# Bottom tip
txt(s, "操作：找一张大白纸，花 1 小时从已有资料和访谈中填满，填不上的格子优先去补齐",
    IN(1.1), IN(5.85), IN(11.33), IN(0.35),
    font=SANS, size=12, color=STONE)
page_num(s, 5, 19)

# ═══════════════════════════════════════════════════════════
# Slide 6: Chapter 2 Divider
# ═══════════════════════════════════════════════════════════

s = blank_slide(BRAND)
chapter_divider(s, "02", "用穿透式问题快速深挖")

# ═══════════════════════════════════════════════════════════
# Slide 7: Dimension 1 - Value & Demand
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第二步 · 穿透式问题", "价值与需求：客户为什么买单，不买会怎样")

tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(IN(1)), Emu(IN(1.9)), Emu(IN(1.1)), Emu(IN(0.4)))
tag.fill.solid(); tag.fill.fore_color.rgb = TAG_BG; tag.line.fill.background(); tag.shadow.inherit = False
txt(s, "维度 1", IN(1), IN(1.92), IN(1.1), IN(0.4),
    font=SERIF, size=13, color=BRAND, align=PP_ALIGN.CENTER)

multi_txt(s, [
    {'text': '1. 客户现在用什么替代方案？那个替代方案哪里不好？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '2. 如果明天涨价50%，客户会骂着接受，还是直接换对手？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '3. 这是"必须满足"的需求（水电级）还是"锦上添花"（盲盒级）？', 'size': 20},
], IN(1), IN(2.55), IN(11.33), IN(3.8))
page_num(s, 7, 19)

# ═══════════════════════════════════════════════════════════
# Slide 8: Dimension 2 - Process & Operations
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第二步 · 穿透式问题", "流程与运转：业务靠哪几个环节活着，哪里最容易断")

tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(IN(1)), Emu(IN(1.9)), Emu(IN(1.1)), Emu(IN(0.4)))
tag.fill.solid(); tag.fill.fore_color.rgb = TAG_BG; tag.line.fill.background(); tag.shadow.inherit = False
txt(s, "维度 2", IN(1), IN(1.92), IN(1.1), IN(0.4),
    font=SERIF, size=13, color=BRAND, align=PP_ALIGN.CENTER)

multi_txt(s, [
    {'text': '1. 关键链路：从客户第一次知道我们到持续使用，必须经过哪5-7个步骤？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '2. 瓶颈环节：当前哪一步最慢/成本最高/最容易出异常？砍掉它业务还成立吗？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '3. 依赖关系：内部必须靠哪个岗位决策？外部必须靠哪个供应商活着？', 'size': 20},
], IN(1), IN(2.55), IN(11.33), IN(3.8))
page_num(s, 8, 19)

# ═══════════════════════════════════════════════════════════
# Slide 9: Dimension 3 - Financials
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第二步 · 穿透式问题", "财务与数字：增长靠什么杠杆，单位经济能不能跑通")

tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(IN(1)), Emu(IN(1.9)), Emu(IN(1.1)), Emu(IN(0.4)))
tag.fill.solid(); tag.fill.fore_color.rgb = TAG_BG; tag.line.fill.background(); tag.shadow.inherit = False
txt(s, "维度 3", IN(1), IN(1.92), IN(1.1), IN(0.4),
    font=SERIF, size=13, color=BRAND, align=PP_ALIGN.CENTER)

multi_txt(s, [
    {'text': '1. 核心杠杆：增长靠"更多客户"还是"更贵客单价"还是"更高复购"？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '2. 单位经济：服务一个客户，收的钱减去直接成本后是正还是负？', 'size': 20},
    {'text': '', 'size': 14},
    {'text': '3. 效率标尺：人均产出、库存周转天数、回款周期——和行业平均比是优是劣？', 'size': 20},
], IN(1), IN(2.55), IN(11.33), IN(3.8))
page_num(s, 9, 19)

# ═══════════════════════════════════════════════════════════
# Slide 10: Dimension 4 - Risk
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第二步 · 穿透式问题", "风险与不确定性：什么事件能让业务3个月内崩溃")

tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(IN(1)), Emu(IN(1.9)), Emu(IN(1.1)), Emu(IN(0.4)))
tag.fill.solid(); tag.fill.fore_color.rgb = TAG_BG; tag.line.fill.background(); tag.shadow.inherit = False
txt(s, "维度 4", IN(1), IN(1.92), IN(1.1), IN(0.4),
    font=SERIF, size=13, color=BRAND, align=PP_ALIGN.CENTER)

multi_txt(s, [
    {'text': '1. 最大的一颗雷：什么事件发生，业务会在3个月内崩溃？', 'size': 20},
    {'text': '    常见答案：核心人员离职、政策禁止、关键技术失效、现金流断裂', 'size': 16, 'color': STONE},
    {'text': '', 'size': 14},
    {'text': '2. 竞对最可能从哪里抄袭并击败我们？价格、体验还是捆绑销售？', 'size': 20},
], IN(1), IN(2.55), IN(11.33), IN(3.8))
page_num(s, 10, 19)

# ═══════════════════════════════════════════════════════════
# Slide 11: Chapter 3 Divider
# ═══════════════════════════════════════════════════════════

s = blank_slide(BRAND)
chapter_divider(s, "03", "三套调研动作交叉验证")

# ═══════════════════════════════════════════════════════════
# Slide 12: Three research methods overview
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第三步 · 调研动作", "别只看文档或听汇报，三件事同时做，互相核对")

methods = [
    ("A", "访谈", "分层问不同角色",
     "业务负责人 → 战略目标\n一线销售 → 客户真实抱怨\n财务运营 → 数据和异常\n外部客户 → 没有你们怎么办"),
    ("B", "看数据", "重点找「反常」",
     "近6个月核心指标趋势\n20%客户是否贡献80%收入\n运营快照：积压和异常\n无数据时用公开材料"),
    ("C", "代入角色", "走一遍完整流程",
     "电商：注册→付款→售后\nB2B：线索→签约→回款\n制造：采购→入库→出库\n记录每一步卡点"),
]

cw = IN(3.5); ch = IN(4.0)
gx = IN(0.4)
for i, (letter, title, subtitle, body) in enumerate(methods):
    x = IN(1.15) + i * (cw + gx)
    y = IN(1.9)
    card(s, x, y, cw, ch)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(x + IN(0.2)), Emu(y + IN(0.18)),
                               Emu(IN(0.48)), Emu(IN(0.48)))
    badge.fill.solid(); badge.fill.fore_color.rgb = BRAND
    badge.line.fill.background(); badge.shadow.inherit = False
    txt(s, letter, Emu(x + IN(0.2)), Emu(y + IN(0.18)),
        IN(0.48), IN(0.48), font=SERIF, size=20, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + IN(0.82), y + IN(0.18), cw - IN(1.0), IN(0.5),
        font=SERIF, size=20, color=NEAR_BLACK)
    txt(s, subtitle, x + IN(0.82), y + IN(0.68), cw - IN(1.0), IN(0.35),
        font=SANS, size=13, color=OLIVE)
    txt(s, body, x + IN(0.2), y + IN(1.2), cw - IN(0.4), IN(2.5),
        font=SANS, size=16, color=DARK_WARM)

page_num(s, 12, 19)

# ═══════════════════════════════════════════════════════════
# Slide 13: Interview detail (上)
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第三步 · 访谈细节", "分层问不同角色，每类人问不同的东西（1/2）")

roles_top = [
    ("业务负责人", "战略目标、最大的假设、最担心的事", "容易讲愿景，忽略细节"),
    ("一线销售/客服", "客户最常抱怨什么？丢单真实原因？", "能听到真问题，但可能情绪化"),
]

for i, (role, ask, trap) in enumerate(roles_top):
    y = IN(1.9) + i * IN(1.65)
    card(s, IN(1.1), y, IN(11.1), IN(1.45))
    label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(IN(1.4)), Emu(y + IN(0.2)),
                               Emu(IN(1.6)), Emu(IN(0.48)))
    label.fill.solid(); label.fill.fore_color.rgb = TAG_BG
    label.line.fill.background(); label.shadow.inherit = False
    txt(s, role, IN(1.4), y + IN(0.2), IN(1.6), IN(0.48),
        font=SANS, size=15, color=BRAND, align=PP_ALIGN.CENTER)
    txt(s, ask, IN(3.3), y + IN(0.15), IN(4.8), IN(1.0),
        font=SANS, size=17, color=DARK_WARM)
    txt(s, "⚠ " + trap, IN(8.5), y + IN(0.15), IN(3.4), IN(1.0),
        font=SANS, size=14, color=STONE)

page_num(s, 13, 19)

# ═══════════════════════════════════════════════════════════
# Slide 14: Interview detail (下)
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第三步 · 访谈细节", "分层问不同角色，每类人问不同的东西（2/2）")

roles_bot = [
    ("财务/运营", "实际数据表现、异常指标、预算分配", "数据可能滞后，要问为什么"),
    ("外部客户", "如果没有你们，会怎么办？最看重什么？", "需要匿名，或假装普通用户"),
]

for i, (role, ask, trap) in enumerate(roles_bot):
    y = IN(1.9) + i * IN(1.65)
    card(s, IN(1.1), y, IN(11.1), IN(1.45))
    label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(IN(1.4)), Emu(y + IN(0.2)),
                               Emu(IN(1.6)), Emu(IN(0.48)))
    label.fill.solid(); label.fill.fore_color.rgb = TAG_BG
    label.line.fill.background(); label.shadow.inherit = False
    txt(s, role, IN(1.4), y + IN(0.2), IN(1.6), IN(0.48),
        font=SANS, size=15, color=BRAND, align=PP_ALIGN.CENTER)
    txt(s, ask, IN(3.3), y + IN(0.15), IN(4.8), IN(1.0),
        font=SANS, size=17, color=DARK_WARM)
    txt(s, "⚠ " + trap, IN(8.5), y + IN(0.15), IN(3.4), IN(1.0),
        font=SANS, size=14, color=STONE)

page_num(s, 14, 19)

# ═══════════════════════════════════════════════════════════
# Slide 15: Chapter 4 Divider
# ═══════════════════════════════════════════════════════════

s = blank_slide(BRAND)
chapter_divider(s, "04", "输出一页纸健康度诊断")

# ═══════════════════════════════════════════════════════════
# Slide 16: Diagnostic template
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "第四步 · 一页纸诊断", "不写长报告，填一张表加三个结论")

dims = ["客户需求", "业务逻辑", "财务状况", "运营效率", "风险敞口"]
y = IN(1.85)
# Header
card(s, IN(1.1), y, IN(11.1), IN(0.48), fill=BRAND, border=BRAND, bw=0)
txt(s, "维度", IN(1.3), y + IN(0.07), IN(2), IN(0.35),
    font=SANS, size=13, color=WHITE)
txt(s, "当前状态（客观描述）", IN(3.3), y + IN(0.07), IN(4.2), IN(0.35),
    font=SANS, size=13, color=WHITE)
txt(s, "健康吗？", IN(7.8), y + IN(0.07), IN(1.5), IN(0.35),
    font=SANS, size=13, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "关键证据或疑问点", IN(9.5), y + IN(0.07), IN(2.5), IN(0.35),
    font=SANS, size=13, color=WHITE)

for i, dim in enumerate(dims):
    y = IN(2.42) + i * IN(0.58)
    bg = IVORY if i % 2 == 0 else PARCHMENT
    card(s, IN(1.1), y, IN(11.1), IN(0.5), fill=bg, border=BORDER, bw=0.3)
    txt(s, dim, IN(1.3), y + IN(0.09), IN(2), IN(0.35),
        font=SANS, size=14, color=DARK_WARM)

# Three conclusions
txt(s, "三个结论：", IN(1.1), IN(5.4), IN(2), IN(0.35),
    font=SERIF, size=17, color=NEAR_BLACK)
conclusions = [
    "1. 这个业务最关键的 3 个成功要素是什么？",
    "2. 如果我去负责，第一个月解决哪个具体问题？",
    "3. 用一句话向陌生人解释清楚这个业务。",
]
for i, c in enumerate(conclusions):
    txt(s, c, IN(1.3), IN(5.8) + i * IN(0.25), IN(10.7), IN(0.25),
        font=SANS, size=13, color=OLIVE)

page_num(s, 16, 19)

# ═══════════════════════════════════════════════════════════
# Slide 17: Appendix - Business types
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "附录 · 快速透视技巧", "不同业务类型，抓住核心问题就够了")

biz_types = [
    ("SaaS", "客户流失率（NDR）？实施成本多高？切换成本高还是低？"),
    ("电商", "流量从哪里来？复购率？退货率？库存占用多少现金？"),
    ("内容/社区", "用户为什么留下来？创作者为什么持续生产？靠什么变现？"),
    ("制造业", "产能利用率？订单账期？原材料价格波动影响多大？"),
    ("线下连锁", "单店盈亏平衡点？回本周期？选址逻辑是什么？"),
]

for i, (name, question) in enumerate(biz_types):
    y = IN(1.9) + i * IN(0.95)
    card(s, IN(1.1), y, IN(11.1), IN(0.82))
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(IN(1.4)), Emu(y + IN(0.15)),
                             Emu(IN(1.4)), Emu(IN(0.48)))
    tag.fill.solid(); tag.fill.fore_color.rgb = TAG_BG
    tag.line.fill.background(); tag.shadow.inherit = False
    txt(s, name, IN(1.4), y + IN(0.15), IN(1.4), IN(0.48),
        font=SERIF, size=16, color=BRAND, align=PP_ALIGN.CENTER)
    txt(s, question, IN(3.1), y + IN(0.12), IN(8.9), IN(0.55),
        font=SANS, size=16, color=DARK_WARM)

page_num(s, 17, 19)

# ═══════════════════════════════════════════════════════════
# Slide 18: Summary
# ═══════════════════════════════════════════════════════════

s = blank_slide()
page_header(s, "总结", "四步法回顾")

summary_items = [
    ("01", "骨架模型", "9个问题装满一张纸，信息缺口一目了然"),
    ("02", "穿透式问题", "四个维度的灵魂拷问，逼迫你进入本质"),
    ("03", "交叉验证", "访谈 + 数据 + 代入角色，三者互相核对"),
    ("04", "一页纸诊断", "不写长报告，一张表 + 三个结论收口"),
]

for i, (num, title, desc) in enumerate(summary_items):
    y = IN(2.0) + i * IN(1.15)
    card(s, IN(1.1), y, IN(11.1), IN(0.95))
    txt(s, num, IN(1.4), y + IN(0.15), IN(0.7), IN(0.55),
        font=SERIF, size=26, color=BRAND)
    txt(s, title, IN(2.4), y + IN(0.1), IN(3), IN(0.45),
        font=SERIF, size=22, color=NEAR_BLACK)
    txt(s, desc, IN(2.4), y + IN(0.55), IN(9.5), IN(0.35),
        font=SANS, size=15, color=OLIVE)

page_num(s, 18, 19)

# ═══════════════════════════════════════════════════════════
# Slide 19: Closing
# ═══════════════════════════════════════════════════════════

s = blank_slide()
txt(s, "谢谢", IN(1), IN(2.8), IN(11.33), IN(1.2),
    font=SERIF, size=48, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
brand_bar(s, IN(6.17), IN(4.2), IN(1))
txt(s, "Q&A", IN(1), IN(4.5), IN(11.33), IN(0.6),
    font=SANS, size=20, color=OLIVE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════

prs.save('/Users/sundanian/Documents/projects/ai-agents/my-agent/workspace/业务快速上手分享/如何快速了解一个陌生业务.pptx')
print("OK: 演示文稿已生成，共 19 页")
